import os
from io import BytesIO
from pathlib import Path

# Try to import libraries, if they don't exist the user must install them first.
from docx import Document
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

FIXTURE_DIR = Path(__file__).parent / "documents"
FIXTURE_DIR.mkdir(parents=True, exist_ok=True)

def generate_digital_pdf():
    path = FIXTURE_DIR / "digital.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawString(100, 750, "Digital PDF Title")
    c.drawString(100, 730, "This is a paragraph in a digital PDF.")
    c.showPage()
    c.drawString(100, 750, "Second Page")
    c.save()

def generate_scanned_pdf():
    # A scanned PDF usually contains just an image per page
    path = FIXTURE_DIR / "scanned.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    # create a dummy image
    img = Image.new('RGB', (200, 100), color=(73, 109, 137))
    d = ImageDraw.Draw(img)
    d.text((10, 10), "Scanned Text", fill=(255, 255, 0))
    img_path = FIXTURE_DIR / "temp.png"
    img.save(img_path)
    c.drawImage(str(img_path), 100, 600, width=200, height=100)
    c.showPage()
    c.save()
    img_path.unlink()

def generate_mixed_pdf():
    path = FIXTURE_DIR / "mixed.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawString(100, 750, "Mixed PDF")
    
    img = Image.new('RGB', (100, 50), color='red')
    img_path = FIXTURE_DIR / "temp2.png"
    img.save(img_path)
    c.drawImage(str(img_path), 100, 650, width=100, height=50)
    c.showPage()
    c.save()
    img_path.unlink()

def generate_docx():
    path = FIXTURE_DIR / "headings_tables.docx"
    doc = Document()
    doc.add_heading('Docx Title', 0)
    doc.add_heading('Section 1', level=1)
    doc.add_paragraph('This is a paragraph.')
    doc.add_paragraph('List item 1', style='List Bullet')
    doc.add_paragraph('List item 2', style='List Bullet')
    
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = 'Header 1'
    table.cell(0, 1).text = 'Header 2'
    table.cell(1, 0).text = 'Row 1 Col 1'
    table.cell(1, 1).text = 'Row 1 Col 2'
    doc.save(str(path))

def generate_pptx():
    path = FIXTURE_DIR / "tables_charts.pptx"
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Presentation Title"
    
    # Table slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Table Slide"
    x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
    table = slide.shapes.add_table(2, 2, x, y, cx, cy).table
    table.cell(0, 0).text = 'Col 1'
    table.cell(0, 1).text = 'Col 2'
    table.cell(1, 0).text = 'Val 1'
    table.cell(1, 1).text = 'Val 2'
    
    # Empty slide
    slide_layout = prs.slide_layouts[6]
    prs.slides.add_slide(slide_layout)
    
    prs.save(str(path))

def generate_xlsx():
    path = FIXTURE_DIR / "multi_sheet.xlsx"
    wb = Workbook()
    
    ws1 = wb.active
    ws1.title = "Sheet 1"
    ws1['A1'] = "Header A"
    ws1['B1'] = "Header B"
    ws1['A2'] = 10
    ws1['B2'] = 20
    ws1['C2'] = "=A2+B2"
    ws1.merge_cells('A3:B3')
    ws1['A3'] = "Merged Cell"
    
    ws2 = wb.create_sheet(title="Sheet 2")
    ws2['A1'] = "Empty sheet test"
    
    wb.save(str(path))

def generate_image():
    path = FIXTURE_DIR / "flowchart.png"
    img = Image.new('RGB', (300, 200), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.rectangle([50, 50, 150, 100], outline="black", width=2)
    d.text((60, 60), "Start", fill="black")
    d.line([150, 75, 200, 75], fill="black", width=2)
    d.rectangle([200, 50, 280, 100], outline="black", width=2)
    d.text((210, 60), "End", fill="black")
    img.save(str(path))

if __name__ == "__main__":
    generate_digital_pdf()
    generate_scanned_pdf()
    generate_mixed_pdf()
    generate_docx()
    generate_pptx()
    generate_xlsx()
    generate_image()
    print("Mock documents generated successfully.")
